from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO
import hashlib
import json
import math
import mimetypes
import os
from pathlib import Path
import re
import sqlite3
from typing import Any, Iterable

import pandas as pd
import pdfplumber
import requests
from fastapi import HTTPException, UploadFile, status

from backend.app.core.config import get_settings
from backend.app.db.connection import Database

try:
    import chromadb
except ImportError:  # pragma: no cover - optional dependency
    chromadb = None

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover - optional dependency
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None

try:
    from sentence_transformers import SentenceTransformer
except ImportError:  # pragma: no cover - optional dependency
    SentenceTransformer = None


RESEARCH_CATEGORIES = [
    "Survey Dataset",
    "Research Report",
    "Final Report",
    "Presentation Deck",
    "Interview Transcript",
    "Focus Group Discussion",
    "Call Summary",
    "Audio Recording",
    "Other Research Documents",
]

AUDIO_EXTENSIONS = {".wav", ".mp3", ".m4a", ".ogg", ".flac", ".aac"}
TEXT_EXTENSIONS = {".txt", ".md", ".rtf", ".json"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".csv", ".xlsx", ".xls"}


@dataclass(slots=True)
class RetrievedChunk:
    document_id: int
    document_name: str
    file_type: str
    category: str
    chunk_index: int
    page_number: int | None
    section: str
    chunk_text: str
    score: float
    metadata: dict[str, Any]


class KnowledgeBaseService:
    def __init__(self) -> None:
        self.settings = get_settings()
        self.database = Database()
        self._embedding_model: Any | None = None
        self._chroma_client: Any | None = None

    def create_client(self, name: str, slug: str | None = None, admin_username: str = "admin", admin_password: str | None = None) -> dict[str, Any]:
        """Create a new client organization with an initial admin user."""
        client_name = name.strip()
        client_slug = self._slugify(slug or client_name)

        if not client_name:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Client name is required")
        if not client_slug:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Client slug could not be generated")

        from backend.app.core.security import hash_password

        admin_password = admin_password or "Password123!"
        admin_hash = hash_password(admin_password)
        admin_email = f"admin@{client_slug}.local"

        with self.database.connect() as connection:
            # Check uniqueness
            existing = connection.execute(
                "SELECT id FROM clients WHERE slug = ?", (client_slug,)
            ).fetchone()
            if existing is not None:
                raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Client with slug '{client_slug}' already exists")

            cursor = connection.execute(
                "INSERT INTO clients (name, slug) VALUES (?, ?)",
                (client_name, client_slug),
            )
            client_id = cursor.lastrowid

            connection.execute(
                """
                INSERT INTO users (client_id, email, username, full_name, password_hash, is_active)
                VALUES (?, ?, ?, ?, ?, 1)
                """,
                (client_id, admin_email, admin_username, f"{client_name} Admin", admin_hash),
            )
            connection.commit()

        return {
            "id": client_id,
            "name": client_name,
            "slug": client_slug,
            "admin_username": admin_username,
            "admin_email": admin_email,
            "admin_password": admin_password,
            "message": f"Client '{client_name}' created. Login with username '{admin_username}' and password '{admin_password}'",
        }

    def list_clients(self) -> list[dict[str, Any]]:
        with self.database.connect() as connection:
            rows = connection.execute("SELECT id, name, slug, created_at FROM clients ORDER BY name").fetchall()
            result = []
            for row in rows:
                project_count = connection.execute(
                    "SELECT COUNT(*) FROM projects WHERE client_id = ?", (row["id"],)
                ).fetchone()[0]
                document_count = connection.execute(
                    "SELECT COUNT(*) FROM documents WHERE client_id = ?", (row["id"],)
                ).fetchone()[0]
                result.append({
                    "id": row["id"],
                    "name": row["name"],
                    "slug": row["slug"],
                    "created_at": row["created_at"],
                    "project_count": project_count,
                    "document_count": document_count,
                })
            return result

    def list_projects(self, client_id: int) -> list[dict[str, Any]]:
        with self.database.connect() as connection:
            projects = connection.execute(
                """
                SELECT id, name, slug, description, is_active
                FROM projects
                WHERE client_id = ?
                ORDER BY name
                """,
                (client_id,),
            ).fetchall()
            result = []
            for project in projects:
                counts = self._project_counts(connection, client_id, project["id"])
                result.append({
                    "id": project["id"],
                    "name": project["name"],
                    "slug": project["slug"],
                    "description": project["description"],
                    "is_active": bool(project["is_active"]),
                    **counts,
                })
            return result

    def create_project(self, client_id: int, name: str, description: str = "", slug: str | None = None) -> dict[str, Any]:
        project_name = name.strip()
        project_description = description.strip()
        project_slug = self._slugify(slug or project_name)

        if not project_name:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Project name is required")
        if not project_slug:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Project slug could not be generated")

        with self.database.connect() as connection:
            unique_slug = project_slug
            suffix = 2
            while connection.execute(
                "SELECT 1 FROM projects WHERE client_id = ? AND slug = ?",
                (client_id, unique_slug),
            ).fetchone() is not None:
                unique_slug = f"{project_slug}-{suffix}"
                suffix += 1

            cursor = connection.execute(
                """
                INSERT INTO projects (client_id, name, slug, description)
                VALUES (?, ?, ?, ?)
                """,
                (client_id, project_name, unique_slug, project_description),
            )
            connection.commit()

        return self.get_project(client_id, cursor.lastrowid)

    def get_project(self, client_id: int, project_id: int) -> dict[str, Any]:
        with self.database.connect() as connection:
            project = connection.execute(
                """
                SELECT p.id, p.name, p.slug, p.description, p.is_active, c.name AS client_name, c.slug AS client_slug
                FROM projects p
                JOIN clients c ON c.id = p.client_id
                WHERE p.id = ? AND p.client_id = ?
                """,
                (project_id, client_id),
            ).fetchone()
            if project is None:
                raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")
            counts = self._project_counts(connection, client_id, project_id)
            return {
                "id": project["id"],
                "name": project["name"],
                "slug": project["slug"],
                "description": project["description"],
                "is_active": bool(project["is_active"]),
                "client_name": project["client_name"],
                "client_slug": project["client_slug"],
                **counts,
            }

    def _slugify(self, value: str) -> str:
        slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower())
        return slug.strip("-")

    def list_documents(self, client_id: int, project_id: int) -> list[dict[str, Any]]:
        with self.database.connect() as connection:
            rows = connection.execute(
                """
                SELECT d.id, d.file_name, d.file_type, d.category, d.status, d.uploaded_at,
                       COALESCE(COUNT(dc.id), 0) AS chunk_count
                FROM documents d
                LEFT JOIN document_chunks dc ON dc.document_id = d.id
                WHERE d.client_id = ? AND d.project_id = ?
                GROUP BY d.id
                ORDER BY d.uploaded_at DESC
                """,
                (client_id, project_id),
            ).fetchall()
            return [
                {
                    "id": row["id"],
                    "file_name": row["file_name"],
                    "file_type": row["file_type"],
                    "category": row["category"],
                    "status": row["status"],
                    "uploaded_at": row["uploaded_at"],
                    "chunk_count": row["chunk_count"],
                }
                for row in rows
            ]

    def project_stats(self, client_id: int, project_id: int) -> dict[str, Any]:
        with self.database.connect() as connection:
            document_count = connection.execute(
                "SELECT COUNT(*) FROM documents WHERE client_id = ? AND project_id = ?",
                (client_id, project_id),
            ).fetchone()[0]
            chunk_count = connection.execute(
                "SELECT COUNT(*) FROM document_chunks WHERE client_id = ? AND project_id = ?",
                (client_id, project_id),
            ).fetchone()[0]
            chat_count = connection.execute(
                "SELECT COUNT(*) FROM chats WHERE client_id = ? AND project_id = ?",
                (client_id, project_id),
            ).fetchone()[0]
            category_rows = connection.execute(
                """
                SELECT category, COUNT(*) AS count
                FROM documents
                WHERE client_id = ? AND project_id = ?
                GROUP BY category
                ORDER BY count DESC, category ASC
                """,
                (client_id, project_id),
            ).fetchall()
            recent_uploads = connection.execute(
                """
                SELECT id, file_name, file_type, category, status, uploaded_at
                FROM documents
                WHERE client_id = ? AND project_id = ?
                ORDER BY uploaded_at DESC
                LIMIT 5
                """,
                (client_id, project_id),
            ).fetchall()
            return {
                "project_id": project_id,
                "document_count": document_count,
                "chunk_count": chunk_count,
                "chat_count": chat_count,
                "category_counts": {row["category"]: row["count"] for row in category_rows},
                "recent_uploads": [dict(row) for row in recent_uploads],
            }

    def ingest_upload(self, client_id: int, project_id: int, upload_file: UploadFile) -> dict[str, Any]:
        project = self.get_project(client_id, project_id)
        raw_bytes = upload_file.file.read()
        if not raw_bytes:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Uploaded file is empty")

        storage_path = self._store_upload(project, upload_file.filename or "upload", raw_bytes)
        extracted_text = self.extract_text(storage_path, upload_file.filename or storage_path.name, upload_file.content_type)
        cleaned_text = self.clean_text(extracted_text)
        if not cleaned_text.strip():
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No readable text could be extracted from the file")

        category = self.classify_document(upload_file.filename or storage_path.name, cleaned_text, upload_file.content_type)
        chunks = self.chunk_text(cleaned_text)
        embeddings = self.embed_texts([chunk["text"] for chunk in chunks])

        with self.database.connect() as connection:
            document_id = self._insert_document(
                connection=connection,
                client_id=client_id,
                project_id=project_id,
                file_name=upload_file.filename or storage_path.name,
                file_type=self.detect_file_type(upload_file.filename or storage_path.name, upload_file.content_type),
                category=category,
                storage_path=str(storage_path),
                status="indexed",
            )
            self._insert_chunks(connection, document_id, client_id, project_id, chunks, embeddings, upload_file.filename or storage_path.name, category)
            connection.commit()

        self._index_to_chroma(project, document_id, chunks, embeddings, upload_file.filename or storage_path.name, category)

        return {
            "document_id": document_id,
            "file_name": upload_file.filename or storage_path.name,
            "file_type": self.detect_file_type(upload_file.filename or storage_path.name, upload_file.content_type),
            "category": category,
            "chunk_count": len(chunks),
            "storage_path": str(storage_path),
            "status": "indexed",
        }

    def ask_question(self, client_id: int, project_id: int, user_id: int, question: str) -> dict[str, Any]:
        project = self.get_project(client_id, project_id)
        rewritten_query = self.rewrite_query(question)
        chunks = self.retrieve_chunks(client_id, project_id, rewritten_query, limit=5)
        answer = self.generate_answer(question=question, rewritten_query=rewritten_query, project=project, chunks=chunks)
        sources = [self._chunk_source(chunk) for chunk in chunks]

        with self.database.connect() as connection:
            connection.execute(
                """
                INSERT INTO chats (client_id, project_id, user_id, question, answer, sources)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (client_id, project_id, user_id, question, answer, json.dumps(sources)),
            )
            connection.commit()

        return {
            "query": rewritten_query,
            "answer": answer,
            "project_id": project_id,
            "document_id": chunks[0].document_id if chunks else None,
            "sources": sources,
        }

    def rewrite_query(self, question: str) -> str:
        normalized = re.sub(r"\s+", " ", question).strip()
        if not normalized:
            return normalized
        if self.settings.groq_api_key:
            prompt = (
                "Rewrite the following research query into a clear, concise search request. "
                "Fix spelling, expand abbreviations, and preserve the original meaning. Return only the rewritten query.\n\n"
                f"Query: {normalized}"
            )
            response_text = self._groq_chat([
                {"role": "system", "content": "You rewrite search queries for a market research assistant."},
                {"role": "user", "content": prompt},
            ])
            if response_text:
                return self._strip_wrappers(response_text)
        replacements = {
            "mumbay": "Mumbai",
            "prcing": "pricing",
            "prcng": "pricing",
            "custmr": "customer",
            "respndent": "respondent",
        }
        for source, target in replacements.items():
            normalized = re.sub(source, target, normalized, flags=re.IGNORECASE)
        return normalized[0].upper() + normalized[1:] if normalized else normalized

    def classify_document(self, file_name: str, text: str, content_type: str | None = None) -> str:
        local_category = self._heuristic_category(file_name, text, content_type)
        if self.settings.groq_api_key:
            prompt = (
                "Classify this research asset into exactly one category from the allowed list. "
                f"Allowed categories: {', '.join(RESEARCH_CATEGORIES)}. "
                "Return only the category name.\n\n"
                f"File name: {file_name}\n\n"
                f"Text sample:\n{text[:4000]}"
            )
            response_text = self._groq_chat([
                {"role": "system", "content": "You classify enterprise research documents."},
                {"role": "user", "content": prompt},
            ])
            if response_text:
                candidate = self._strip_wrappers(response_text)
                if candidate in RESEARCH_CATEGORIES:
                    return candidate
        return local_category

    def detect_file_type(self, file_name: str, content_type: str | None = None) -> str:
        suffix = Path(file_name).suffix.lower()
        if suffix in AUDIO_EXTENSIONS:
            return "audio"
        if suffix == ".pdf":
            return "pdf"
        if suffix == ".docx":
            return "docx"
        if suffix == ".pptx":
            return "pptx"
        if suffix in {".csv", ".xlsx", ".xls"}:
            return "dataset"
        if suffix in TEXT_EXTENSIONS:
            return "text"
        if content_type:
            return content_type.split("/")[0]
        return suffix.lstrip(".") or "unknown"

    def extract_text(self, file_path: Path, file_name: str, content_type: str | None = None) -> str:
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return self._extract_pdf(file_path)
        if suffix == ".docx":
            return self._extract_docx(file_path)
        if suffix == ".pptx":
            return self._extract_pptx(file_path)
        if suffix == ".csv":
            return self._extract_csv(file_path)
        if suffix in {".xlsx", ".xls"}:
            return self._extract_excel(file_path)
        if suffix in TEXT_EXTENSIONS:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        if suffix in AUDIO_EXTENSIONS:
            return self._transcribe_audio(file_path)
        return file_path.read_text(encoding="utf-8", errors="ignore")

    def clean_text(self, text: str) -> str:
        cleaned = text.replace("\x00", " ")
        cleaned = re.sub(r"\r\n?", "\n", cleaned)
        cleaned = re.sub(r"[ \t]+", " ", cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        return cleaned.strip()

    def chunk_text(self, text: str, chunk_size: int = 180, overlap: int = 30) -> list[dict[str, Any]]:
        words = text.split()
        if not words:
            return []
        if len(words) <= chunk_size:
            return [{"chunk_index": 0, "text": text.strip(), "section": "", "page_number": None}]

        step = max(1, chunk_size - overlap)
        chunks: list[dict[str, Any]] = []
        for start in range(0, len(words), step):
            window = words[start : start + chunk_size]
            if not window:
                continue
            chunks.append({
                "chunk_index": len(chunks),
                "text": " ".join(window).strip(),
                "section": self._derive_section(window),
                "page_number": None,
            })
            if start + chunk_size >= len(words):
                break
        return chunks

    def embed_texts(self, texts: Iterable[str]) -> list[list[float]]:
        text_list = list(texts)
        if not text_list:
            return []
        if SentenceTransformer is not None:
            model = self._get_embedding_model()
            if model is not None:
                try:
                    vectors = model.encode(text_list, normalize_embeddings=True)
                    return [vector.tolist() for vector in vectors]
                except Exception:
                    pass
        return [self._hash_embedding(text) for text in text_list]

    def retrieve_chunks(self, client_id: int, project_id: int, query: str, limit: int = 5) -> list[RetrievedChunk]:
        project = self.get_project(client_id, project_id)
        chroma_chunks = self._retrieve_from_chroma(project, query, limit)
        if chroma_chunks:
            return chroma_chunks
        return self._retrieve_from_sqlite(client_id, project_id, query, limit)

    def generate_answer(
        self,
        question: str,
        rewritten_query: str,
        project: dict[str, Any],
        chunks: list[RetrievedChunk],
    ) -> str:
        if not chunks:
            return (
                "I could not find supporting evidence in the selected project for that question. "
                "Try uploading more source material or rephrasing the query."
            )

        context_block = self._format_context_block(chunks)
        if self.settings.groq_api_key:
            prompt = (
                "Answer the user's question using only the provided project context. "
                "Be concise, factual, and include the source document names inline when relevant. "
                "If the context is insufficient, say so directly.\n\n"
                f"Project: {project['name']}\n"
                f"Original question: {question}\n"
                f"Rewritten question: {rewritten_query}\n\n"
                f"Context:\n{context_block}"
            )
            response_text = self._groq_chat([
                {"role": "system", "content": "You are an enterprise research assistant grounded in project evidence."},
                {"role": "user", "content": prompt},
            ])
            if response_text:
                return self._strip_wrappers(response_text)

        source_lines = []
        for chunk in chunks[:3]:
            snippet = self._trim_sentence(chunk.chunk_text, 280)
            source_lines.append(
                f"{chunk.document_name} ({chunk.category}{f', page {chunk.page_number}' if chunk.page_number else ''}): {snippet}"
            )
        return (
            f"Based on the selected project, I found these supporting excerpts for '{question}':\n"
            + "\n".join(f"- {line}" for line in source_lines)
        )

    def rerank_chunks(self, query: str, chunks: list[RetrievedChunk], limit: int | None = None) -> list[RetrievedChunk]:
        if not chunks:
            return []

        query_terms = set(re.findall(r"[A-Za-z0-9]+", query.lower()))
        reranked: list[RetrievedChunk] = []
        for chunk in chunks:
            chunk_terms = set(re.findall(r"[A-Za-z0-9]+", chunk.chunk_text.lower()))
            overlap = len(query_terms & chunk_terms) / max(1, len(query_terms))
            combined_score = (0.65 * float(chunk.score)) + (0.35 * overlap)
            reranked.append(
                RetrievedChunk(
                    document_id=chunk.document_id,
                    document_name=chunk.document_name,
                    file_type=chunk.file_type,
                    category=chunk.category,
                    chunk_index=chunk.chunk_index,
                    page_number=chunk.page_number,
                    section=chunk.section,
                    chunk_text=chunk.chunk_text,
                    score=combined_score,
                    metadata=chunk.metadata,
                )
            )

        reranked.sort(key=lambda item: item.score, reverse=True)
        return reranked[:limit] if limit is not None else reranked

    def verify_answer(
        self,
        question: str,
        rewritten_query: str,
        answer: str,
        chunks: list[RetrievedChunk],
    ) -> tuple[str, str, str]:
        if not chunks:
            fallback = (
                "I could not find sufficient supporting evidence in the selected project to answer that question. "
                "Try uploading more source material or rephrasing the query."
            )
            return fallback, "rejected", "No supporting chunks were retrieved."

        support_terms = set()
        for chunk in chunks:
            support_terms.update(re.findall(r"[A-Za-z0-9]+", chunk.chunk_text.lower()))

        answer_terms = set(re.findall(r"[A-Za-z0-9]+", answer.lower()))
        if not answer_terms:
            return answer, "accepted", "Answer is empty but treated as supported by retrieval."

        coverage = len(answer_terms & support_terms) / max(1, len(answer_terms))
        if coverage < 0.15 and "could not find" not in answer.lower() and "not enough" not in answer.lower():
            fallback = (
                "I could not verify that answer against the retrieved project evidence. "
                "The selected project does not appear to contain enough support for a grounded response."
            )
            return fallback, "rejected", f"Low evidence coverage ({coverage:.2f})."

        if self.settings.groq_api_key:
            verification_prompt = (
                "Review the draft answer against the project evidence. If it is grounded, return the answer unchanged. "
                "If it is unsupported, return a short fallback message stating that evidence is insufficient. "
                "Do not add new facts.\n\n"
                f"Original question: {question}\n"
                f"Rewritten query: {rewritten_query}\n"
                f"Draft answer: {answer}\n\n"
                f"Evidence:\n{self._format_context_block(chunks)}"
            )
            response_text = self._groq_chat([
                {"role": "system", "content": "You verify whether answers are grounded in retrieved evidence."},
                {"role": "user", "content": verification_prompt},
            ])
            if response_text:
                verified = self._strip_wrappers(response_text)
                if verified:
                    status_label = "accepted" if "could not find" not in verified.lower() and "insufficient" not in verified.lower() else "rejected"
                    notes = "Verified with Groq grounding check." if status_label == "accepted" else "Groq grounding check rejected the answer."
                    return verified, status_label, notes

        return answer, "accepted", f"Evidence coverage {coverage:.2f}."

    def save_chat_response(
        self,
        client_id: int,
        project_id: int,
        user_id: int,
        question: str,
        answer: str,
        sources: list[dict[str, Any]],
    ) -> None:
        with self.database.connect() as connection:
            connection.execute(
                """
                INSERT INTO chats (client_id, project_id, user_id, question, answer, sources)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (client_id, project_id, user_id, question, answer, json.dumps(sources)),
            )
            connection.commit()

    def log_retrieval_event(
        self,
        client_id: int,
        project_id: int,
        user_id: int,
        original_query: str,
        rewritten_query: str,
        top_document_id: int | None,
        top_document_name: str,
        top_chunk_score: float,
        verification_status: str,
    ) -> None:
        with self.database.connect() as connection:
            connection.execute(
                """
                INSERT INTO retrieval_events (
                    client_id, project_id, user_id, original_query, rewritten_query,
                    top_document_id, top_document_name, top_chunk_score, verification_status
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    client_id,
                    project_id,
                    user_id,
                    original_query,
                    rewritten_query,
                    top_document_id,
                    top_document_name,
                    top_chunk_score,
                    verification_status,
                ),
            )
            connection.commit()

    def dashboard_overview(self, client_id: int) -> dict[str, Any]:
        with self.database.connect() as connection:
            total_clients = connection.execute(
                "SELECT COUNT(*) FROM clients WHERE id = ?",
                (client_id,),
            ).fetchone()[0]
            total_projects = connection.execute(
                "SELECT COUNT(*) FROM projects WHERE client_id = ?",
                (client_id,),
            ).fetchone()[0]
            total_documents = connection.execute(
                "SELECT COUNT(*) FROM documents WHERE client_id = ?",
                (client_id,),
            ).fetchone()[0]
            total_chunks = connection.execute(
                "SELECT COUNT(*) FROM document_chunks WHERE client_id = ?",
                (client_id,),
            ).fetchone()[0]
            total_chats = connection.execute(
                "SELECT COUNT(*) FROM chats WHERE client_id = ?",
                (client_id,),
            ).fetchone()[0]
            total_retrievals = connection.execute(
                "SELECT COUNT(*) FROM retrieval_events WHERE client_id = ?",
                (client_id,),
            ).fetchone()[0]

            recent_uploads = connection.execute(
                """
                SELECT d.file_name, d.file_type, d.category, d.status, d.uploaded_at, p.name AS project_name
                FROM documents d
                JOIN projects p ON p.id = d.project_id
                WHERE d.client_id = ?
                ORDER BY d.uploaded_at DESC
                LIMIT 10
                """,
                (client_id,),
            ).fetchall()
            recent_queries = connection.execute(
                """
                SELECT original_query, rewritten_query, top_document_name, top_chunk_score, verification_status, created_at
                FROM retrieval_events
                WHERE client_id = ?
                ORDER BY created_at DESC
                LIMIT 10
                """,
                (client_id,),
            ).fetchall()
            project_rows = connection.execute(
                """
                SELECT p.id, p.name, p.slug, p.description, p.is_active,
                       COUNT(DISTINCT d.id) AS document_count,
                       COUNT(DISTINCT dc.id) AS chunk_count,
                       COUNT(DISTINCT c.id) AS chat_count
                FROM projects p
                LEFT JOIN documents d ON d.project_id = p.id
                LEFT JOIN document_chunks dc ON dc.project_id = p.id
                LEFT JOIN chats c ON c.project_id = p.id
                WHERE p.client_id = ?
                GROUP BY p.id
                ORDER BY p.name
                """,
                (client_id,),
            ).fetchall()

            categories = connection.execute(
                """
                SELECT category, COUNT(*) AS count
                FROM documents
                WHERE client_id = ?
                GROUP BY category
                ORDER BY count DESC, category ASC
                """,
                (client_id,),
            ).fetchall()

        return {
            "total_clients": total_clients,
            "total_projects": total_projects,
            "total_documents": total_documents,
            "total_chunks": total_chunks,
            "total_chats": total_chats,
            "total_retrievals": total_retrievals,
            "recent_uploads": [dict(row) for row in recent_uploads],
            "recent_queries": [dict(row) for row in recent_queries],
            "projects": [dict(row) for row in project_rows],
            "category_counts": {row["category"]: row["count"] for row in categories},
        }

    def recent_chat_history(self, client_id: int, project_id: int, limit: int = 20) -> list[dict[str, Any]]:
        with self.database.connect() as connection:
            rows = connection.execute(
                """
                SELECT id, question, answer, sources, created_at
                FROM chats
                WHERE client_id = ? AND project_id = ?
                ORDER BY created_at DESC
                LIMIT ?
                """,
                (client_id, project_id, limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def export_chat_bundle(self, client_id: int, project_id: int, export_kind: str = "chat") -> tuple[bytes, str, str]:
        project = self.get_project(client_id, project_id)
        chat_rows = self.recent_chat_history(client_id, project_id, limit=500)
        stats = self.project_stats(client_id, project_id)
        documents = self.list_documents(client_id, project_id)

        chat_df = pd.DataFrame(chat_rows)
        docs_df = pd.DataFrame(documents)
        summary_df = pd.DataFrame([
            {
                "project_name": project["name"],
                "project_slug": project["slug"],
                "document_count": stats["document_count"],
                "chunk_count": stats["chunk_count"],
                "chat_count": stats["chat_count"],
                "category_counts": json.dumps(stats["category_counts"]),
            }
        ])

        export_kind = export_kind.lower().strip()
        if export_kind not in {"chat", "summary", "context"}:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Unsupported export kind")

        if export_kind == "chat":
            payload = self._dataframe_to_excel(
                {
                    "chat_history": chat_df,
                    "documents": docs_df,
                }
            )
            filename = f"{project['slug']}_chat_export.xlsx"
        elif export_kind == "summary":
            payload = self._dataframe_to_excel(
                {
                    "project_summary": summary_df,
                    "recent_uploads": pd.DataFrame(stats.get("recent_uploads", [])),
                    "category_counts": pd.DataFrame(
                        [{"category": category, "count": count} for category, count in stats.get("category_counts", {}).items()]
                    ),
                }
            )
            filename = f"{project['slug']}_summary_export.xlsx"
        else:
            retrieval_rows = self.recent_retrieval_events(client_id, project_id, limit=100)
            payload = self._dataframe_to_excel(
                {
                    "project_context": summary_df,
                    "documents": docs_df,
                    "retrieval_events": pd.DataFrame(retrieval_rows),
                    "chat_history": chat_df,
                }
            )
            filename = f"{project['slug']}_context_export.xlsx"

        return payload, filename, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

    def delete_project(self, client_id: int, project_id: int) -> dict[str, Any]:
        """Delete a project, its documents, chunks, chroma collection, and uploaded files."""
        project = self.get_project(client_id, project_id)

        # 1. Delete ChromaDB collection
        try:
            if chromadb is not None:
                if self._chroma_client is None:
                    self._chroma_client = chromadb.PersistentClient(path=str(self.settings.chroma_path))
                collection_name = self._collection_name(project)
                try:
                    self._chroma_client.delete_collection(collection_name)
                except Exception:
                    pass  # Collection may not exist
        except Exception:
            pass

        # 2. Delete uploaded files from disk
        project_root = self.settings.upload_path / project["client_slug"] / project["slug"]
        if project_root.exists():
            import shutil
            try:
                shutil.rmtree(project_root)
            except Exception:
                pass

        # 3. Delete from database (cascading via schema)
        with self.database.connect() as connection:
            connection.execute("DELETE FROM document_chunks WHERE client_id = ? AND project_id = ?", (client_id, project_id))
            connection.execute("DELETE FROM documents WHERE client_id = ? AND project_id = ?", (client_id, project_id))
            connection.execute("DELETE FROM chats WHERE client_id = ? AND project_id = ?", (client_id, project_id))
            connection.execute("DELETE FROM retrieval_events WHERE client_id = ? AND project_id = ?", (client_id, project_id))
            connection.execute("DELETE FROM projects WHERE id = ? AND client_id = ?", (project_id, client_id))
            connection.commit()

        return {"deleted": True, "project_id": project_id, "project_name": project["name"]}

    def delete_client(self, client_id: int) -> dict[str, Any]:
        """Delete a client and ALL associated data across all projects."""
        with self.database.connect() as connection:
            client = connection.execute(
                "SELECT id, name, slug FROM clients WHERE id = ?", (client_id,)
            ).fetchone()
            if client is None:
                raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Client not found")

            # Get all projects for ChromaDB cleanup
            projects = connection.execute(
                "SELECT slug FROM projects WHERE client_id = ?", (client_id,)
            ).fetchall()

        # 1. Delete ChromaDB collections for all projects
        try:
            if chromadb is not None:
                if self._chroma_client is None:
                    self._chroma_client = chromadb.PersistentClient(path=str(self.settings.chroma_path))
                client_slug = client["slug"]
                for project_row in projects:
                    project_slug = project_row["slug"]
                    collection_name = f"{client_slug}__{project_slug}"
                    try:
                        self._chroma_client.delete_collection(collection_name)
                    except Exception:
                        pass
        except Exception:
            pass

        # 2. Delete uploaded files from disk
        client_upload_root = self.settings.upload_path / client["slug"]
        if client_upload_root.exists():
            import shutil
            try:
                shutil.rmtree(client_upload_root)
            except Exception:
                pass

        # 3. Delete from database (cascading via schema)
        with self.database.connect() as connection:
            connection.execute("DELETE FROM document_chunks WHERE client_id = ?", (client_id,))
            connection.execute("DELETE FROM documents WHERE client_id = ?", (client_id,))
            connection.execute("DELETE FROM chats WHERE client_id = ?", (client_id,))
            connection.execute("DELETE FROM retrieval_events WHERE client_id = ?", (client_id,))
            connection.execute("DELETE FROM sessions WHERE client_id = ?", (client_id,))
            connection.execute("DELETE FROM users WHERE client_id = ?", (client_id,))
            connection.execute("DELETE FROM projects WHERE client_id = ?", (client_id,))
            connection.execute("DELETE FROM clients WHERE id = ?", (client_id,))
            connection.commit()

        return {"deleted": True, "client_id": client_id, "client_name": client["name"]}

    def delete_document(self, client_id: int, project_id: int, document_id: int) -> dict[str, Any]:
        """Delete a single document, its chunks, uploaded file, and chroma entries."""
        with self.database.connect() as connection:
            doc = connection.execute(
                "SELECT id, file_name, storage_path, file_type FROM documents WHERE id = ? AND client_id = ? AND project_id = ?",
                (document_id, client_id, project_id),
            ).fetchone()
            if doc is None:
                raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")

            # Delete from SQLite
            connection.execute("DELETE FROM document_chunks WHERE document_id = ?", (document_id,))
            connection.execute("DELETE FROM documents WHERE id = ?", (document_id,))
            connection.commit()

        # Delete uploaded file from disk
        if doc["storage_path"]:
            stored_path = Path(doc["storage_path"])
            if stored_path.exists():
                try:
                    stored_path.unlink()
                except Exception:
                    pass

        # Delete from ChromaDB - get all chunk IDs for this document
        project = self.get_project(client_id, project_id)
        try:
            if chromadb is not None:
                collection = self._get_chroma_collection(project)
                # Query for all entries with this document_id
                existing = collection.get(where={"document_id": str(document_id)})
                if existing and existing.get("ids"):
                    collection.delete(ids=existing["ids"])
        except Exception:
            pass

        return {"deleted": True, "document_id": document_id, "file_name": doc["file_name"]}

    def recent_retrieval_events(self, client_id: int, project_id: int, limit: int = 20) -> list[dict[str, Any]]:
        with self.database.connect() as connection:
            rows = connection.execute(
                """
                SELECT original_query, rewritten_query, top_document_id, top_document_name, top_chunk_score,
                       verification_status, created_at
                FROM retrieval_events
                WHERE client_id = ? AND project_id = ?
                ORDER BY created_at DESC
                LIMIT ?
                """,
                (client_id, project_id, limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def _dataframe_to_excel(self, sheet_map: dict[str, pd.DataFrame]) -> bytes:
        buffer = BytesIO()
        with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
            for sheet_name, dataframe in sheet_map.items():
                safe_dataframe = dataframe.copy()
                if safe_dataframe.empty:
                    safe_dataframe = pd.DataFrame([{"note": "No rows available"}])
                safe_dataframe.to_excel(writer, sheet_name=sheet_name[:31], index=False)
        buffer.seek(0)
        return buffer.getvalue()

    def _project_counts(self, connection: sqlite3.Connection, client_id: int, project_id: int) -> dict[str, int]:
        document_count = connection.execute(
            "SELECT COUNT(*) FROM documents WHERE client_id = ? AND project_id = ?",
            (client_id, project_id),
        ).fetchone()[0]
        chunk_count = connection.execute(
            "SELECT COUNT(*) FROM document_chunks WHERE client_id = ? AND project_id = ?",
            (client_id, project_id),
        ).fetchone()[0]
        chat_count = connection.execute(
            "SELECT COUNT(*) FROM chats WHERE client_id = ? AND project_id = ?",
            (client_id, project_id),
        ).fetchone()[0]
        return {"document_count": document_count, "chunk_count": chunk_count, "chat_count": chat_count}

    def _store_upload(self, project: dict[str, Any], file_name: str, raw_bytes: bytes) -> Path:
        project_root = self.settings.upload_path / project["client_slug"] / project["slug"]
        project_root.mkdir(parents=True, exist_ok=True)
        safe_name = re.sub(r"[^A-Za-z0-9._-]+", "_", Path(file_name).name)
        timestamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
        storage_path = project_root / f"{timestamp}_{safe_name}"
        storage_path.write_bytes(raw_bytes)
        return storage_path

    def _insert_document(
        self,
        connection: sqlite3.Connection,
        client_id: int,
        project_id: int,
        file_name: str,
        file_type: str,
        category: str,
        storage_path: str,
        status: str,
    ) -> int:
        cursor = connection.execute(
            """
            INSERT INTO documents (client_id, project_id, file_name, file_type, category, storage_path, status)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (client_id, project_id, file_name, file_type, category, storage_path, status),
        )
        return cursor.lastrowid

    def _insert_chunks(
        self,
        connection: sqlite3.Connection,
        document_id: int,
        client_id: int,
        project_id: int,
        chunks: list[dict[str, Any]],
        embeddings: list[list[float]],
        file_name: str,
        category: str,
    ) -> None:
        for chunk, embedding in zip(chunks, embeddings, strict=False):
            metadata = {
                "client_id": client_id,
                "project_id": project_id,
                "document_id": document_id,
                "document_name": file_name,
                "category": category,
                "chunk_index": chunk["chunk_index"],
                "section": chunk.get("section", ""),
                "page_number": chunk.get("page_number"),
            }
            connection.execute(
                """
                INSERT INTO document_chunks (
                    document_id, client_id, project_id, chunk_index, section, page_number,
                    chunk_text, metadata_json, embedding_json
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    document_id,
                    client_id,
                    project_id,
                    chunk["chunk_index"],
                    chunk.get("section", ""),
                    chunk.get("page_number"),
                    chunk["text"],
                    json.dumps(metadata),
                    json.dumps(embedding),
                ),
            )

    def _index_to_chroma(
        self,
        project: dict[str, Any],
        document_id: int,
        chunks: list[dict[str, Any]],
        embeddings: list[list[float]],
        file_name: str,
        category: str,
    ) -> None:
        if chromadb is None or not chunks:
            return
        try:
            collection = self._get_chroma_collection(project)
            ids = [f"{document_id}-{chunk['chunk_index']}" for chunk in chunks]
            metadatas = [
                {
                    "client_id": str(project.get("client_id", "")),
                    "project_id": str(project["id"]),
                    "document_id": str(document_id),
                    "document_name": file_name,
                    "category": category,
                    "chunk_index": chunk["chunk_index"],
                    "section": chunk.get("section", ""),
                    "page_number": chunk.get("page_number") or -1,
                }
                for chunk in chunks
            ]
            collection.upsert(
                ids=ids,
                documents=[chunk["text"] for chunk in chunks],
                embeddings=embeddings,
                metadatas=metadatas,
            )
        except Exception:
            return

    def _retrieve_from_chroma(self, project: dict[str, Any], query: str, limit: int) -> list[RetrievedChunk]:
        if chromadb is None:
            return []
        try:
            collection = self._get_chroma_collection(project)
            query_embedding = self.embed_texts([query])[0]
            result = collection.query(
                query_embeddings=[query_embedding],
                n_results=limit,
                include=["documents", "metadatas", "distances"],
            )
            chunks: list[RetrievedChunk] = []
            documents = result.get("documents", [[]])[0]
            metadatas = result.get("metadatas", [[]])[0]
            distances = result.get("distances", [[]])[0]
            for index, (document_text, metadata, distance) in enumerate(zip(documents, metadatas, distances, strict=False)):
                chunks.append(
                    RetrievedChunk(
                        document_id=int(metadata.get("document_id", 0)),
                        document_name=metadata.get("document_name", ""),
                        file_type=metadata.get("file_type", ""),
                        category=metadata.get("category", ""),
                        chunk_index=int(metadata.get("chunk_index", index)),
                        page_number=metadata.get("page_number") if metadata.get("page_number", -1) != -1 else None,
                        section=metadata.get("section", ""),
                        chunk_text=document_text,
                        score=max(0.0, 1.0 - float(distance)),
                        metadata=dict(metadata),
                    )
                )
            return chunks
        except Exception:
            return []

    def _retrieve_from_sqlite(self, client_id: int, project_id: int, query: str, limit: int) -> list[RetrievedChunk]:
        with self.database.connect() as connection:
            rows = connection.execute(
                """
                SELECT dc.document_id, dc.chunk_index, dc.section, dc.page_number, dc.chunk_text,
                       dc.metadata_json, dc.embedding_json,
                       d.file_name, d.file_type, d.category
                FROM document_chunks dc
                JOIN documents d ON d.id = dc.document_id
                WHERE dc.client_id = ? AND dc.project_id = ?
                """,
                (client_id, project_id),
            ).fetchall()
            if not rows:
                return []
            query_embedding = self.embed_texts([query])[0]
            scored: list[RetrievedChunk] = []
            for row in rows:
                stored_embedding = json.loads(row["embedding_json"])
                score = self._cosine_similarity(query_embedding, stored_embedding)
                scored.append(
                    RetrievedChunk(
                        document_id=row["document_id"],
                        document_name=row["file_name"],
                        file_type=row["file_type"],
                        category=row["category"],
                        chunk_index=row["chunk_index"],
                        page_number=row["page_number"],
                        section=row["section"],
                        chunk_text=row["chunk_text"],
                        score=score,
                        metadata=json.loads(row["metadata_json"]),
                    )
                )
            scored.sort(key=lambda item: item.score, reverse=True)
            return scored[:limit]

    def _get_chroma_collection(self, project: dict[str, Any]) -> Any:
        if chromadb is None:
            raise RuntimeError("ChromaDB is not available")
        if self._chroma_client is None:
            self._chroma_client = chromadb.PersistentClient(path=str(self.settings.chroma_path))
        collection_name = self._collection_name(project)
        return self._chroma_client.get_or_create_collection(name=collection_name, metadata={"hnsw:space": "cosine"})

    def _collection_name(self, project: dict[str, Any]) -> str:
        client_slug = self._slugify(project.get("client_slug", str(project.get("client_id", "client"))))
        project_slug = self._slugify(project.get("slug", str(project["id"])))
        return f"{client_slug}__{project_slug}"

    def _get_embedding_model(self) -> Any | None:
        if self._embedding_model is not None:
            return self._embedding_model
        if SentenceTransformer is None:
            return None
        try:
            self._embedding_model = SentenceTransformer(self.settings.embedding_model)
        except Exception:
            self._embedding_model = None
        return self._embedding_model

    def _heuristic_category(self, file_name: str, text: str, content_type: str | None = None) -> str:
        filename = file_name.lower()
        sample = f"{file_name}\n{text[:2000]}".lower()
        if Path(file_name).suffix.lower() in AUDIO_EXTENSIONS:
            return "Audio Recording"
        if any(keyword in filename for keyword in ["survey", "dataset", "csv", "xls", "xlsx"]):
            return "Survey Dataset"
        if any(keyword in sample for keyword in ["focus group", "fgd", "group discussion"]):
            return "Focus Group Discussion"
        if any(keyword in sample for keyword in ["interview", "transcript", "verbatim", "respondent"]):
            return "Interview Transcript"
        if any(keyword in sample for keyword in ["call summary", "customer call", "sales call"]):
            return "Call Summary"
        if any(keyword in sample for keyword in ["presentation", "slides", "slide", "deck"]):
            return "Presentation Deck"
        if any(keyword in sample for keyword in ["final report", "final deck"]):
            return "Final Report"
        if any(keyword in sample for keyword in ["research report", "insights report", "analysis report"]):
            return "Research Report"
        return "Other Research Documents"

    def _transcribe_audio(self, file_path: Path) -> str:
        try:
            import whisper  # type: ignore
        except ImportError as exc:  # pragma: no cover - optional dependency
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Audio transcription requires the openai-whisper package to be installed",
            ) from exc
        try:
            model = whisper.load_model("base")
            result = model.transcribe(str(file_path))
            return result.get("text", "")
        except Exception as exc:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Audio transcription failed") from exc

    def _extract_pdf(self, file_path: Path) -> str:
        texts: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                texts.append(page.extract_text() or "")
        return "\n\n".join(texts)

    def _extract_docx(self, file_path: Path) -> str:
        if DocxDocument is None:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="python-docx is required to read Word documents")
        document = DocxDocument(str(file_path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n".join(paragraphs)

    def _extract_pptx(self, file_path: Path) -> str:
        if Presentation is None:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="python-pptx is required to read PowerPoint files")
        presentation = Presentation(str(file_path))
        slide_text: list[str] = []
        for slide in presentation.slides:
            collected: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    collected.append(shape.text.strip())
            if collected:
                slide_text.append("\n".join(collected))
        return "\n\n".join(slide_text)

    def _extract_csv(self, file_path: Path) -> str:
        dataframe = pd.read_csv(file_path)
        return dataframe.fillna("").astype(str).to_string(index=False)

    def _extract_excel(self, file_path: Path) -> str:
        sheets = pd.read_excel(file_path, sheet_name=None)
        blocks: list[str] = []
        for sheet_name, dataframe in sheets.items():
            sheet_text = dataframe.fillna("").astype(str).to_string(index=False)
            blocks.append(f"Sheet: {sheet_name}\n{sheet_text}")
        return "\n\n".join(blocks)

    def _groq_chat(self, messages: list[dict[str, str]]) -> str:
        if not self.settings.groq_api_key:
            return ""
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {self.settings.groq_api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": self.settings.groq_model,
                "messages": messages,
                "temperature": 0.1,
            },
            timeout=45,
        )
        response.raise_for_status()
        payload = response.json()
        choices = payload.get("choices", [])
        if not choices:
            return ""
        return choices[0]["message"]["content"].strip()

    def _hash_embedding(self, text: str, dimensions: int = 384) -> list[float]:
        vector = [0.0] * dimensions
        tokens = re.findall(r"[A-Za-z0-9]+", text.lower())
        if not tokens:
            return vector
        for token in tokens:
            digest = hashlib.sha256(token.encode("utf-8")).hexdigest()
            index = int(digest[:8], 16) % dimensions
            weight = 1.0 + (int(digest[8:12], 16) % 7) / 10.0
            vector[index] += weight
        norm = math.sqrt(sum(value * value for value in vector)) or 1.0
        return [value / norm for value in vector]

    def _cosine_similarity(self, left: list[float], right: list[float]) -> float:
        if not left or not right:
            return 0.0
        length = min(len(left), len(right))
        dot_product = sum(left[index] * right[index] for index in range(length))
        left_norm = math.sqrt(sum(value * value for value in left[:length])) or 1.0
        right_norm = math.sqrt(sum(value * value for value in right[:length])) or 1.0
        return dot_product / (left_norm * right_norm)

    def _chunk_source(self, chunk: RetrievedChunk) -> dict[str, Any]:
        source = {
            "document_id": chunk.document_id,
            "document_name": chunk.document_name,
            "file_type": chunk.file_type,
            "category": chunk.category,
            "chunk_index": chunk.chunk_index,
            "page_number": chunk.page_number,
            "section": chunk.section,
            "score": round(chunk.score, 4),
            "snippet": self._trim_sentence(chunk.chunk_text, 240),
        }
        return source

    def _format_context_block(self, chunks: list[RetrievedChunk]) -> str:
        parts: list[str] = []
        for chunk in chunks:
            label = f"{chunk.document_name} | {chunk.category}"
            if chunk.page_number is not None:
                label += f" | page {chunk.page_number}"
            if chunk.section:
                label += f" | section {chunk.section}"
            parts.append(f"[{label}]\n{chunk.chunk_text}")
        return "\n\n".join(parts)

    def _trim_sentence(self, text: str, length: int) -> str:
        cleaned = re.sub(r"\s+", " ", text).strip()
        if len(cleaned) <= length:
            return cleaned
        trimmed = cleaned[:length].rsplit(" ", 1)[0]
        return f"{trimmed}..."

    def _derive_section(self, words: list[str]) -> str:
        preview = " ".join(words[:10]).strip()
        return preview[:80]

    def _strip_wrappers(self, text: str) -> str:
        stripped = text.strip()
        stripped = stripped.removeprefix("\"").removesuffix("\"")
        stripped = stripped.removeprefix("'").removesuffix("'")
        return stripped

    def _slugify(self, value: str) -> str:
        slug = re.sub(r"[^A-Za-z0-9]+", "-", value).strip("-").lower()
        return slug or "project"
