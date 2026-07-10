"""Improved document extraction with better structure preservation."""
from __future__ import annotations

import re
from pathlib import Path
import pandas as pd

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None


def extract_pptx_improved(file_path: Path) -> str:
    """Extract text from PowerPoint with explicit slide labeling."""
    if Presentation is None:
        raise RuntimeError("python-pptx is required to read PowerPoint files")
    
    presentation = Presentation(str(file_path))
    slide_texts: list[str] = []
    
    for slide_num, slide in enumerate(presentation.slides, start=1):
        collected: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Detect if this looks like a title (short, uppercase, or bold formatting hint)
                if len(text) < 100 and (text.isupper() or slide_num == 1):
                    collected.append(f"SLIDE {slide_num} TITLE: {text}")
                else:
                    collected.append(text)
        if collected:
            slide_texts.append(f"=== SLIDE {slide_num} ===\n" + "\n".join(collected))
    
    return "\n\n".join(slide_texts)


def extract_excel_improved(file_path: Path) -> str:
    """Extract Excel with explicit sheet labeling and table structure."""
    sheets = pd.read_excel(file_path, sheet_name=None)
    blocks: list[str] = []
    
    for sheet_num, (sheet_name, df) in enumerate(sheets.items(), start=1):
        # Clean the dataframe
        df_clean = df.fillna("").astype(str)
        
        # Detect if first row looks like headers
        headers = df_clean.iloc[0].tolist() if len(df_clean) > 0 else []
        has_headers = any(h and re.match(r"^[A-Z_][a-zA-Z0-9_\s]*$", str(h)) for h in headers)
        
        # Format as structured text for better retrieval
        if has_headers and len(headers) > 2:
            # Format as table: Header1 | Header2 | Header3
            table_lines = []
            header_row = " | ".join(str(h).strip() for h in headers)
            table_lines.append(f"Sheet: {sheet_name}")
            table_lines.append(f"Columns: {header_row}")
            table_lines.append("-" * len(header_row))
            
            for _, row in df_clean.iterrows():
                row_values = [str(v).strip()[:30] for v in row.tolist()]
                table_lines.append(" | ".join(row_values))
            
            blocks.append("\n".join(table_lines))
        else:
            # Regular format without headers
            sheet_text = df_clean.to_string(index=False)
            blocks.append(f"Sheet: {sheet_name}\n{sheet_text}")
    
    return "\n\n".join(blocks)


def extract_csv_improved(file_path: Path) -> str:
    """Extract CSV with column headers highlighted."""
    df = pd.read_csv(file_path)
    df_clean = df.fillna("").astype(str)
    
    header_row = " | ".join(df_clean.columns.tolist())
    
    lines = []
    lines.append(f"Columns: {header_row}")
    lines.append("-" * len(header_row))
    
    for _, row in df_clean.iterrows():
        row_values = [str(v).strip()[:30] for v in row.tolist()]
        lines.append(" | ".join(row_values))
    
    return "\n".join(lines)


def extract_docx_improved(file_path: Path) -> str:
    """Extract Word doc with heading preservation."""
    if DocxDocument is None:
        raise RuntimeError("python-docx is required to read Word documents")
    
    document = DocxDocument(str(file_path))
    
    # Try to get headings from styles
    structured_lines: list[str] = []
    
    for para in document.paragraphs:
        if not para.text.strip():
            continue
        
        # Check if it's a heading (based on style name)
        style_name = para.style.name.lower() if para.style else ""
        if "heading" in style_name or "title" in style_name:
            structured_lines.append(f"HEADING: {para.text.strip()}")
        elif "list" in style_name:
            structured_lines.append(f"LIST ITEM: {para.text.strip()}")
        else:
            structured_lines.append(para.text.strip())
    
    return "\n".join(structured_lines)


if __name__ == "__main__":
    # Test the extractors
    import sys
    if len(sys.argv) < 2:
        print("Usage: python extractors.py <file_path>")
        sys.exit(1)
    
    file_path = Path(sys.argv[1])
    suffix = file_path.suffix.lower()
    
    if suffix == ".pptx":
        print(extract_pptx_improved(file_path))
    elif suffix in {".xlsx", ".xls"}:
        print(extract_excel_improved(file_path))
    elif suffix == ".csv":
        print(extract_csv_improved(file_path))
    elif suffix == ".docx":
        print(extract_docx_improved(file_path))
    else:
        print(f"Unsupported file type: {suffix}")